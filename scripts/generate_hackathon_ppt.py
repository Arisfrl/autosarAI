from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BASE_DIR = Path(__file__).resolve().parent.parent
OUTPUT_PATH = BASE_DIR / "AUTOSAR_AI_MDE_Hackathon_Demo.pptx"
SCREENSHOT_DIR = BASE_DIR / "screenshots"
STREAMLIT_SCREENSHOT = SCREENSHOT_DIR / "streamlit_ui_rich.png"
STREAMLIT_COMPILE_SCREENSHOT = SCREENSHOT_DIR / "streamlit_compile_safety.png"
HMI_SCREENSHOT = SCREENSHOT_DIR / "hmi_cluster_ui.png"

BG = RGBColor(9, 21, 35)
SURFACE = RGBColor(16, 34, 56)
SURFACE_2 = RGBColor(19, 43, 73)
ACCENT = RGBColor(45, 226, 168)
ACCENT_2 = RGBColor(57, 184, 255)
TEXT = RGBColor(232, 243, 255)
MUTED = RGBColor(184, 208, 229)
WARN = RGBColor(255, 191, 84)
DANGER = RGBColor(255, 98, 98)
WHITE = RGBColor(255, 255, 255)

SLIDES = [
    {
        "layout": "title",
        "title": "AUTOSAR AI MDE Hackathon Project",
        "subtitle": "AI-Augmented Model-Driven Engineering for Software-Defined Vehicles\nNatural language to YAML to Adaptive + Classic ARXML with safety checks and demo HMI",
    },
    {
        "title": "Problem We Are Solving",
        "bullets": [
            "AUTOSAR design is slow, artifact-heavy, and split across Classic CAN and Adaptive SOME/IP stacks.",
            "Teams must turn informal requirements into models, mappings, and safety evidence under time pressure.",
            "Early iterations often lack fast validation, traceability, and reusable demo artifacts.",
            "We need a working pipeline that converts intent and ARXML references into usable outputs, not just documents.",
        ],
        "callouts": [
            ("Pain", "Manual modeling, fragmented communication stacks, late safety feedback", DANGER),
            ("Need", "Fast, local, explainable AUTOSAR artifact generation", ACCENT),
        ],
    },
    {
        "title": "What This Project Delivers",
        "bullets": [
            "Natural language to AUTOSAR YAML using local grounded retrieval.",
            "One-click compile to both Adaptive ARXML and Classic ARXML.",
            "Safety validation for structure, ASIL coverage, mappings, and communication risk.",
            "Reference ARXML upload, AI suggestions, HMI simulation, and tenant-scoped audit trails.",
        ],
        "callouts": [
            ("Artifacts", "YAML, Adaptive ARXML, Classic ARXML, safety findings", ACCENT_2),
            ("UX", "Streamlit app with guided demo buttons and downloads", ACCENT),
        ],
    },
    {
        "title": "End-to-End Workflow",
        "bullets": [
            "1. User provides a natural-language AUTOSAR request or loads an AEB / CAN-to-SOME-IP demo prompt.",
            "2. Engine retrieves relevant local references from markdown and AUTOSAR example ARXML documents.",
            "3. Local Ollama or Gemini generates normalized YAML for system, ECUs, services, signals, and safety.",
            "4. Core compiler converts that YAML into normalized AUTOSAR-like XML outputs for both Adaptive and Classic variants.",
            "5. Safety validator and mapping baseline score the result before users export artifacts or drive the HMI simulation.",
        ],
        "callouts": [
            ("Core pipeline", "Prompt -> RAG -> YAML -> ARXML -> Safety/HMI", WARN),
        ],
    },
    {
        "title": "How AI Is Helping",
        "bullets": [
            "RAG grounds prompts in local AUTOSAR notes and example ARXML, reducing generic LLM drift.",
            "The model converts broad intent into structured system, ECU, service, signal, and safety sections.",
            "AI suggestions propose safety, mapping, and deployment improvements after generation.",
            "A lightweight ML baseline flags risky signal-to-service mappings before compile time.",
        ],
        "callouts": [
            ("LLM", "Ollama default, Gemini optional provider", ACCENT_2),
            ("ML", "Mapping baseline test accuracy 0.80, F1 0.8571", ACCENT),
        ],
    },
    {
        "title": "What Is In The Repo",
        "bullets": [
            "core.py: ingestion, retrieval, YAML parsing, ARXML compilation, safety validation, AI suggestions, and mapping pre-check logic.",
            "streamlit_app.py: authenticated multi-tenant UI, demo flows, model selection, upload validation, compile/download actions, and safety reporting.",
            "AUTOSAR_WorkflowExample/ and data/: local AUTOSAR reference material used for grounding and example architectures.",
            "scripts/ and models/: baseline ML training/inference utilities and saved mapping metrics.",
            "audit/ and chroma_db/: tenant-scoped audit trails and retrieval persistence for BMW/Ford demo isolation.",
        ],
        "callouts": [
            ("Security", "Login gate + tenant-scoped storage + audit JSONL", WARN),
            ("Demo assets", "SVG visuals, HMI cluster HTML, ARXML simulator", ACCENT_2),
        ],
    },
    {
        "title": "Safety, Governance, and Validation",
        "bullets": [
            "Safety report checks system naming, ECU ASIL coverage, hardware bindings, service presence, signal mappings, and safety section completeness.",
            "Mapping pre-check can block compile when the ML baseline sees a high-confidence invalid or unsafe mapping.",
            "ARXML uploads are classified as Classic / Adaptive / Unknown / Invalid before being accepted as references.",
            "Audit logs track generation, compile, reload, and fallback actions per tenant, supporting demo governance and traceability.",
        ],
        "callouts": [
            ("Validation", "Structure + mapping + safety + upload sanity checks", DANGER),
        ],
    },
    {
        "title": "Current Evidence From The Project",
        "bullets": [
            "Mapping baseline metrics: train accuracy 0.9474, validation accuracy 1.00, test accuracy 0.80, test F1 0.8571.",
            "Current end-to-end batch result snapshot: 1 pass / 10 cases with most failures caused by local Ollama timeout, not by parser or compiler logic.",
            "The app includes deterministic fallback YAML so the demo remains usable even when model inference is unreliable.",
            "Recent enhancement: one compile action now emits both Adaptive and Classic ARXML outputs for the same generated architecture.",
        ],
        "callouts": [
            ("Strength", "Architecture toolchain works end-to-end", ACCENT),
            ("Gap", "LLM latency/reliability is the current bottleneck", DANGER),
        ],
    },
    {
        "title": "Live Demo Story",
        "bullets": [
            "Log in as BMW or Ford demo admin to show tenant isolation and auditability.",
            "Load the AEB or CAN-to-SOME-IP prompt, generate YAML, and explain how retrieval grounds the output.",
            "Compile once and show both Adaptive and Classic ARXML outputs plus downloadable files.",
            "Run the safety check and AI suggestion flow to show how the tool catches issues and proposes improvements.",
            "Open the HMI cluster HTML to show ARXML-driven telltale authorization and human-facing validation behavior.",
        ],
        "callouts": [
            ("Best demo path", "Prompt -> YAML -> dual ARXML -> safety -> HMI", ACCENT_2),
        ],
    },
    {
        "layout": "screenshots",
        "title": "Demo Screenshots",
        "subtitle": "Authenticated Streamlit console and ARXML-driven HMI cluster",
    },
    {
        "layout": "compile_safety",
        "title": "Compile and Safety Evidence",
        "subtitle": "Dual ARXML output and safety validation results captured from the live app",
    },
    {
        "title": "Why This Matters",
        "bullets": [
            "Faster concept-to-architecture turnaround for SDV teams working across Classic and Adaptive AUTOSAR stacks.",
            "Earlier safety feedback reduces rework and makes ASIL-oriented design tradeoffs visible sooner.",
            "Local-first inference and tenant separation make the approach suitable for sensitive automotive design environments.",
            "The same framework can evolve from hackathon scaffold into a stronger engineering assistant with richer schema validation and enterprise identity.",
        ],
        "callouts": [
            ("Benefit", "Speed, consistency, safety visibility, and demo-grade explainability", ACCENT),
        ],
    },
    {
        "title": "Judging Summary",
        "bullets": [
            "Innovation: combines local LLM generation, RAG, safety checks, ML mapping pre-check, and HMI simulation in one AUTOSAR demo flow.",
            "Technical depth: covers YAML generation, dual ARXML compilation, tenant isolation, validation, audit logging, and UI-driven review.",
            "Practical impact: accelerates AUTOSAR architecture iteration while exposing safety and communication issues earlier.",
            "Demo strength: judges can see prompt input, generated artifacts, compile evidence, safety results, and HMI behavior in one session.",
        ],
        "callouts": [
            ("One-line pitch", "An AI-assisted AUTOSAR copilot that turns intent into dual-stack artifacts with safety-aware feedback.", ACCENT_2),
            ("Judge lens", "Innovation + feasibility + impact + live proof", WARN),
        ],
    },
    {
        "title": "Next Steps",
        "bullets": [
            "Replace simplified XML with richer AUTOSAR package/schema support and stronger Classic vs Adaptive specialization.",
            "Improve LLM robustness and latency through prompt slimming, model tuning, caching, or better local runtime setup.",
            "Expand ML datasets and features for communication mapping, hazard patterns, and safety mechanism recommendations.",
            "Add official AUTOSAR document ingestion, OAuth/OIDC, and richer reporting for production-ready engineering workflows.",
        ],
        "callouts": [
            ("Hackathon message", "We already proved the flow; next step is industrial-grade depth", WARN),
        ],
    },
]


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_top_bar(slide, title_text):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = SURFACE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title_text
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Pt(18)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_footer(slide, text="AUTOSAR AI MDE Hackathon Demo"):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.05), Inches(5.5), Inches(0.28))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED


def add_bullet_box(slide, bullets, x, y, w, h, font_size=20):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for index, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)
        p.bullet = True
    return box


def add_callout(slide, title, body, color, x, y, w, h):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE_2
    shape.line.color.rgb = color
    shape.line.width = Pt(1.8)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = "Aptos Display"
    r1.font.size = Pt(16)
    r1.font.bold = True
    r1.font.color.rgb = color
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Aptos"
    r2.font.size = Pt(14)
    r2.font.color.rgb = TEXT


def add_metric_row(slide, items):
    left = Inches(0.55)
    top = Inches(5.55)
    width = Inches(12.2)
    gap = Inches(0.18)
    card_w = (width - gap * (len(items) - 1)) / len(items)
    for idx, (label, value, color) in enumerate(items):
        x = left + idx * (card_w + gap)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, top, card_w, Inches(1.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE
        shape.line.color.rgb = color
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = value
        r1.font.name = "Aptos Display"
        r1.font.size = Pt(20)
        r1.font.bold = True
        r1.font.color.rgb = color
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = label
        r2.font.name = "Aptos"
        r2.font.size = Pt(11)
        r2.font.color.rgb = MUTED


def add_image_with_frame(slide, image_path, x, y, w, h, caption):
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    frame.fill.solid()
    frame.fill.fore_color.rgb = SURFACE
    frame.line.color.rgb = ACCENT_2
    frame.line.width = Pt(1.5)

    if image_path.exists():
        slide.shapes.add_picture(str(image_path), x + Inches(0.08), y + Inches(0.08), w - Inches(0.16), h - Inches(0.48))

    cap = slide.shapes.add_textbox(x + Inches(0.1), y + h - Inches(0.33), w - Inches(0.2), Inches(0.2))
    tf = cap.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = caption
    r.font.name = "Aptos"
    r.font.size = Pt(11)
    r.font.color.rgb = MUTED


def add_small_callout(slide, title, body, color, x, y, w, h):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(10, 21, 34)
    shape.line.color.rgb = color
    shape.line.width = Pt(1.3)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = "Aptos Display"
    r1.font.size = Pt(12)
    r1.font.bold = True
    r1.font.color.rgb = color
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Aptos"
    r2.font.size = Pt(10)
    r2.font.color.rgb = TEXT


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for idx, spec in enumerate(SLIDES):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        set_background(slide)

        if spec.get("layout") == "title":
            title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(10.8), Inches(1.5))
            tf = title_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = spec["title"]
            r.font.name = "Aptos Display"
            r.font.size = Pt(30)
            r.font.bold = True
            r.font.color.rgb = WHITE

            sub = slide.shapes.add_textbox(Inches(0.7), Inches(2.25), Inches(8.7), Inches(1.4))
            tf2 = sub.text_frame
            tf2.clear()
            p2 = tf2.paragraphs[0]
            r2 = p2.add_run()
            r2.text = spec["subtitle"]
            r2.font.name = "Aptos"
            r2.font.size = Pt(18)
            r2.font.color.rgb = MUTED

            accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(4.45), Inches(3.0), Inches(1.2))
            accent.fill.solid()
            accent.fill.fore_color.rgb = SURFACE_2
            accent.line.color.rgb = ACCENT
            tf3 = accent.text_frame
            tf3.clear()
            p3 = tf3.paragraphs[0]
            p3.alignment = PP_ALIGN.CENTER
            r3 = p3.add_run()
            r3.text = "Demo focus\nClassic + Adaptive + Safety"
            r3.font.name = "Aptos Display"
            r3.font.size = Pt(16)
            r3.font.bold = True
            r3.font.color.rgb = ACCENT

            add_metric_row(slide, [
                ("Tenants", "BMW + Ford", ACCENT_2),
                ("Outputs", "YAML + 2x ARXML", ACCENT),
                ("Checks", "Safety + Mapping", WARN),
            ])
            add_footer(slide)
            continue

        if spec.get("layout") == "screenshots":
            add_top_bar(slide, spec["title"])

            sub = slide.shapes.add_textbox(Inches(0.65), Inches(0.82), Inches(8.5), Inches(0.35))
            tf = sub.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = spec["subtitle"]
            r.font.name = "Aptos"
            r.font.size = Pt(16)
            r.font.color.rgb = MUTED

            add_image_with_frame(
                slide,
                STREAMLIT_SCREENSHOT,
                Inches(0.65),
                Inches(1.3),
                Inches(7.0),
                Inches(5.25),
                "Streamlit console: authenticated multi-tenant UI, YAML generation, compile, and safety workflow",
            )
            add_image_with_frame(
                slide,
                HMI_SCREENSHOT,
                Inches(7.9),
                Inches(1.3),
                Inches(4.75),
                Inches(5.25),
                "ARXML-driven HMI cluster: telltale authorization, control panel, and validation feedback",
            )
            add_small_callout(slide, "1. Prompt to YAML", "Natural-language request generates structured AUTOSAR YAML in-app.", ACCENT, Inches(0.95), Inches(5.95), Inches(2.2), Inches(0.72))
            add_small_callout(slide, "2. HMI linkage", "Loaded ARXML authorizes telltales before the cluster can be controlled.", WARN, Inches(8.25), Inches(5.95), Inches(3.7), Inches(0.72))
            add_footer(slide)
            continue

        if spec.get("layout") == "compile_safety":
            add_top_bar(slide, spec["title"])

            sub = slide.shapes.add_textbox(Inches(0.65), Inches(0.82), Inches(8.8), Inches(0.35))
            tf = sub.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = spec["subtitle"]
            r.font.name = "Aptos"
            r.font.size = Pt(16)
            r.font.color.rgb = MUTED

            add_image_with_frame(
                slide,
                STREAMLIT_COMPILE_SCREENSHOT,
                Inches(0.65),
                Inches(1.2),
                Inches(9.1),
                Inches(5.75),
                "Live app state showing Classic ARXML output and safety status after compile/check",
            )
            add_small_callout(slide, "Dual output", "One compile flow now generates both Adaptive and Classic ARXML downloads.", ACCENT_2, Inches(10.0), Inches(1.55), Inches(2.65), Inches(0.95))
            add_small_callout(slide, "Safety feedback", "Warnings surface immediately so the demo shows engineering governance, not just generation.", DANGER, Inches(10.0), Inches(2.8), Inches(2.65), Inches(1.1))
            add_small_callout(slide, "Judge takeaway", "This is a working pipeline: model output, compile artifact, and validation evidence in one UI.", ACCENT, Inches(10.0), Inches(4.25), Inches(2.65), Inches(1.15))
            add_footer(slide)
            continue

        add_top_bar(slide, spec["title"])
        add_bullet_box(slide, spec["bullets"], Inches(0.65), Inches(0.95), Inches(7.7), Inches(5.7), font_size=18)

        callouts = spec.get("callouts", [])
        if callouts:
            x = Inches(8.55)
            y = Inches(1.05)
            for title, body, color in callouts:
                add_callout(slide, title, body, color, x, y, Inches(4.05), Inches(1.35))
                y += Inches(1.55)

        if idx == 7:
            add_metric_row(slide, [
                ("Test accuracy", "0.80", ACCENT),
                ("Test F1", "0.8571", ACCENT_2),
                ("E2E pass count", "1 / 10", WARN),
            ])
        elif idx == 8:
            add_metric_row(slide, [
                ("Prompt", "AEB / CAN->SOME-IP", ACCENT),
                ("Compile", "Adaptive + Classic", ACCENT_2),
                ("Review", "Safety + HMI", WARN),
            ])
        else:
            add_metric_row(slide, [
                ("Local-first", "Ollama default", ACCENT),
                ("RAG", "Chroma + references", ACCENT_2),
                ("Governance", "Audit + tenant scope", WARN),
            ])

        add_footer(slide)

    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


if __name__ == "__main__":
    path = build_deck()
    print(f"Created presentation: {path}")
